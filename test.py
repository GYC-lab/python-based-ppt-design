from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from utils import *
from themeColor import *

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_shape_transparency(shape, alpha):
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))

# Create presentation
prs = Presentation()

# Add a slide (empty slide layout)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add a blue box to the slide
blueBox = slide.shapes.add_shape(autoshape_type_id=MSO_SHAPE.RECTANGLE,
                                 left=Cm(0),
                                 top=Cm(0),
                                 height=Cm(10),
                                 width=Cm(20))

# Make the box blue
blueBoxFill = blueBox.fill
blueBoxFill.solid()
blueBoxFillColour = blueBoxFill.fore_color
blueBoxFillColour.rgb = RGBColor(0,176,240)

# Set the transparency of the blue box to 56%
set_solid_transparency(blueBoxFill, 0.8)

myThemeColor  = myThemeColor_user
NewThemeColor = set_my_theme_color(myThemeColor)
set_fill_gradient_color(blueBoxFill,NewThemeColor,2,90)


# Save the presentation
prs.save('./output/motiongo科技1_new.pptx')
