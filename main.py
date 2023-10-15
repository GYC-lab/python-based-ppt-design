from pptx import Presentation
from utils import *
from themeColor import *
from setColor import *

# --------------------------------- settings --------------------------------- #
_change_font_color       = 1 # 1: change font color; 0: not change font color
_change_outline_color    = 0 # 1: change outline color; 0: not change outline color
_change_background_color = 1 # 1: change background color; 0: not change background color
_change_shape_color      = 1 # 1: change shape color; 0: not change fill color
_gradient_angle          = 90 # angle of gradient, 90: vertical; 0: horizontal
_transparency            = 0  # transparency of fill color, 0: not transparent; 1: transparent
_set_single_slile        = 0 # 1: only set one slide; 0: set all slides
slide_index              = 0 # index of slide to be set (if set_single_slile=1)

# --------------------------------- filepath --------------------------------- #
input_ppt_path  = './input/'
output_ppt_path = './output/'
inFile          = 'motiongo科技1_sub1.pptx' # 目录页
inFile          = 'motiongo科技1_sub2.pptx' # 科技对商业的影响 - 背景介绍
inFile          = 'motiongo科技1_sub4.pptx' # 科技发展现状大纲
inFile          = 'motiongo科技1_sub5.pptx' # 标题页
inFile          = 'motiongo科技1_sub3.pptx' # 科技对商业的影响 - 趋势分析
outFile         = 'motiongo科技1_new.pptx'
inFile          = 'motiongo科技1.pptx'      # 全部

# --------------------- set theme color at themeColor.py --------------------- #
myThemeColor  = myThemeColor_user
NewThemeColor = set_my_theme_color(myThemeColor)

# --------------------------------- load ppt --------------------------------- #
prs = Presentation(input_ppt_path+inFile)  
print('loading ppt from %s' % input_ppt_path+inFile)

# ----------------------------------- main ----------------------------------- #
report_choice(_change_font_color,_change_outline_color,_change_background_color,_change_shape_color)
for slide in prs.slides:
    if _set_single_slile:
        if slide_index == prs.slides.index(slide):
            set_color(slide,NewThemeColor,_gradient_angle,_transparency,_change_font_color,_change_outline_color,_change_background_color,_change_shape_color)
            break
    else:
        set_color(slide,NewThemeColor,_gradient_angle,_transparency,_change_font_color,_change_outline_color,_change_background_color,_change_shape_color)

# --------------------------------- save ppt --------------------------------- #
print('done! saving ppt to %s' % output_ppt_path+outFile)
prs.save(output_ppt_path+outFile)
